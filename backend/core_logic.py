import os
import img2pdf
from PIL import Image
import qrcode
from pdf2docx import Converter
from pptx import Presentation
import fitz

def convert_image_to_pdf(image_path: str, output_path: str) -> str:
    # Open image to check validity
    with Image.open(image_path) as img:
        # Convert to RGB if it's RGBA (PNG with transparency) to avoid img2pdf errors
        if img.mode in ("RGBA", "LA", "P"):
            rgb_img = img.convert("RGB")
            temp_img_path = image_path + ".jpg"
            rgb_img.save(temp_img_path, "JPEG")
            try:
                with open(output_path, "wb") as f:
                    f.write(img2pdf.convert(temp_img_path))
            finally:
                if os.path.exists(temp_img_path):
                    os.remove(temp_img_path)
        else:
            with open(output_path, "wb") as f:
                f.write(img2pdf.convert(image_path))
    return output_path

def compress_pdf_file(input_path: str, output_path: str, compression_level: str = "minimum") -> str:
    # Use PyMuPDF (fitz) for PDF compression
    doc = fitz.open(input_path)
    if compression_level == "maximum":
        # Maximum compression removes unreferenced objects, linearizes, and deflates content
        doc.save(output_path, garbage=4, deflate=True, clean=True, linear=True)
    else:
        # Minimum/Recommended compression to simply deflate streams
        doc.save(output_path, garbage=1, deflate=True)
    doc.close()
        
    return output_path

def create_qr_code(text: str, output_path: str) -> str:
    qr = qrcode.QRCode(
        version=1,
        error_correction=qrcode.constants.ERROR_CORRECT_L,
        box_size=10,
        border=4,
    )
    qr.add_data(text)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")
    img.save(output_path)
    return output_path

def convert_pdf_to_docx(pdf_path: str, docx_path: str) -> str:
    # Using pdf2docx
    cv = Converter(pdf_path)
    cv.convert(docx_path, start=0, end=None)
    cv.close()
    return docx_path

def extract_text_from_ppt(ppt_path: str, text_path: str) -> str:
    prs = Presentation(ppt_path)
    text_runs = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
                    
    with open(text_path, "w", encoding="utf-8") as f:
        f.write("\n".join(text_runs))
        
    return text_path
